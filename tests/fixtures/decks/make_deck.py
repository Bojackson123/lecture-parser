"""Generate the mock lecture deck fixtures: lecture01.pdf, lecture01.pptx, value_iteration.png.

Standalone — imports nothing from ``lecturenotes``. Run with the ``fixtures`` group:

    uv run --group fixtures python tests/fixtures/decks/make_deck.py

Deterministic by construction: fixed document metadata, reportlab's ``invariant`` mode,
and PPTX zip entries rewritten with a fixed timestamp, so regenerating without a content
change produces a byte-identical diff. Slide content is ASCII only because reportlab's
built-in fonts are Latin-1; the same strings go into both files so Phase 2 and Phase 4
tests can compare across formats.
"""

from __future__ import annotations

import io
import zipfile
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

HERE = Path(__file__).resolve().parent
PDF_PATH = HERE / "lecture01.pdf"
PPTX_PATH = HERE / "lecture01.pptx"
PNG_PATH = HERE / "value_iteration.png"

AUTHOR = "lecturenotes fixtures"
TITLE = "Lecture 1: MDPs and the Bellman Equation"
FIXED_DATETIME = datetime(2026, 1, 5, 9, 0, 0)
ZIP_DATE_TIME = (2026, 1, 5, 9, 0, 0)

# --- Slide content ---------------------------------------------------------------------

SLIDE1_TITLE = "Markov Decision Processes"
SLIDE1_BULLETS = [
    "States s in S: everything the agent needs to know about the world right now",
    "Actions a in A: the choices available in each state",
    "Reward R(s, a): a number received for taking action a in state s",
    "Transition function T(s, a, s'): where you land after each action",
    "Discount factor gamma in [0, 1): how much later rewards count",
]

SLIDE2_TITLE = "The Bellman Equation"
SLIDE2_LEFT = [
    "Equation",
    "V(s) = max_a [ R(s, a) + gamma * sum_s' T(s, a, s') V(s') ]",
    "V(s): value of state s",
    "max_a: best action available",
    "sum_s' T(s, a, s'): expectation over next states",
    "gamma: discount factor",
]
SLIDE2_RIGHT = [
    "Intuition",
    "Value = immediate reward + discounted future",
    "Recursive: V appears on both sides",
    "Exactly one fixed point",
    "Solve it by repeated substitution",
    "Everything else in the course builds on this",
]

SLIDE3_TITLE = "Value Iteration"
SLIDE3_STEPS = [
    "1. Initialise V_0(s) = 0 for every state",
    "2. Sweep: V_{k+1}(s) = max_a [ R(s, a) +",
    "       gamma * sum_s' T(s, a, s') V_k(s') ]",
    "3. Stop when max_s |V_{k+1}(s) - V_k(s)| < epsilon",
    "4. Read off the greedy policy pi(s) = argmax_a [ ... ]",
]

NOTES = [
    "Start with the four ingredients and keep coming back to them. Students usually "
    "confuse the reward function with the value function, so stress that the reward is "
    "immediate and given. The discount factor is introduced here only so the equation on "
    "the next slide makes sense.",
    "Walk the left column top to bottom before touching the intuition column. The "
    "recursion is the point: the same V appears on both sides. Mention that this will be "
    "on the exam.",
    "Run through one sweep by hand if time allows. The plot is the maximum change between "
    "successive sweeps on a small grid world; it falls geometrically with rate gamma. "
    "Close by previewing model-free methods next week.",
]

# --- Figure ----------------------------------------------------------------------------


def make_png() -> bytes:
    """A tiny convergence plot: max change per sweep, decaying geometrically."""
    w, h = 240, 150
    img = Image.new("P", (w, h), 0)
    img.putpalette([255, 255, 255, 0, 0, 0, 30, 90, 200] + [0, 0, 0] * 253)
    d = ImageDraw.Draw(img)
    x0, y0, x1, y1 = 30, 12, 228, 122
    d.line([(x0, y0), (x0, y1), (x1, y1)], fill=1)
    points = []
    for k in range(12):
        x = x0 + int((x1 - x0) * k / 11)
        y = y1 - int((y1 - y0) * 0.9 * (0.7**k))
        points.append((x, y))
    d.line(points, fill=2, width=2)
    for x, y in points:
        d.rectangle([x - 2, y - 2, x + 2, y + 2], fill=2)
    d.text((x0 + 60, y1 + 8), "sweep k", fill=1)
    d.text((2, y0 - 4), "max", fill=1)
    d.text((2, y0 + 8), "chg", fill=1)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


# --- PDF --------------------------------------------------------------------------------


def make_pdf(png: bytes) -> None:
    page_w, page_h = landscape(A4)
    c = canvas.Canvas(str(PDF_PATH), pagesize=(page_w, page_h), invariant=1, pageCompression=0)
    c.setAuthor(AUTHOR)
    c.setTitle(TITLE)
    c.setSubject("lecturenotes test fixture")
    c.setCreator("tests/fixtures/decks/make_deck.py")

    def title(text: str) -> None:
        c.setFont("Helvetica-Bold", 28)
        c.drawString(60, page_h - 80, text)

    def footer(n: int) -> None:
        c.setFont("Helvetica", 10)
        c.drawRightString(page_w - 60, 30, f"Lecture 1 - slide {n} / 3")

    # Slide 1: single column.
    title(SLIDE1_TITLE)
    c.setFont("Helvetica", 16)
    for i, line in enumerate(SLIDE1_BULLETS):
        c.drawString(80, page_h - 150 - 40 * i, f"- {line}")
    footer(1)
    c.showPage()

    # Slide 2: two columns, drawn row by row so both columns share y-coordinates and the
    # content stream interleaves them. Naive extraction reads L1 R1 L2 R2 ...
    title(SLIDE2_TITLE)
    left_x, right_x = 60, 450
    for i, (left, right) in enumerate(zip(SLIDE2_LEFT, SLIDE2_RIGHT, strict=True)):
        y = page_h - 150 - 40 * i
        c.setFont("Helvetica-Bold" if i == 0 else "Helvetica", 13)
        c.drawString(left_x, y, left)
        c.drawString(right_x, y, right)
    footer(2)
    c.showPage()

    # Slide 3: single column plus a figure on the right.
    title(SLIDE3_TITLE)
    c.setFont("Helvetica", 14)
    for i, line in enumerate(SLIDE3_STEPS):
        c.drawString(80, page_h - 150 - 34 * i, line)
    c.drawImage(ImageReader(io.BytesIO(png)), page_w - 320, 140, width=240, height=150)
    footer(3)
    c.showPage()
    c.save()


# --- PPTX -------------------------------------------------------------------------------


def _fill(placeholder: object, lines: list[str]) -> None:
    tf = placeholder.text_frame  # type: ignore[attr-defined]
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(20)


def make_pptx(png: bytes) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_and_content = prs.slide_layouts[1]
    two_content = prs.slide_layouts[3]

    s1 = prs.slides.add_slide(title_and_content)
    s1.shapes.title.text = SLIDE1_TITLE
    _fill(s1.placeholders[1], SLIDE1_BULLETS)

    s2 = prs.slides.add_slide(two_content)
    s2.shapes.title.text = SLIDE2_TITLE
    _fill(s2.placeholders[1], SLIDE2_LEFT)
    _fill(s2.placeholders[2], SLIDE2_RIGHT)

    s3 = prs.slides.add_slide(title_and_content)
    s3.shapes.title.text = SLIDE3_TITLE
    body = s3.placeholders[1]
    body.width = Inches(7.8)
    _fill(body, SLIDE3_STEPS)
    s3.shapes.add_picture(io.BytesIO(png), Inches(8.6), Inches(2.6), width=Inches(4.0))

    for slide, note in zip(prs.slides, NOTES, strict=True):
        slide.notes_slide.notes_text_frame.text = note

    cp = prs.core_properties
    cp.author = AUTHOR
    cp.last_modified_by = AUTHOR
    cp.title = TITLE
    cp.subject = "lecturenotes test fixture"
    cp.created = FIXED_DATETIME
    cp.modified = FIXED_DATETIME
    cp.last_printed = FIXED_DATETIME
    cp.revision = 1

    raw = io.BytesIO()
    prs.save(raw)
    _rewrite_zip_with_fixed_timestamps(raw.getvalue(), PPTX_PATH)


def _rewrite_zip_with_fixed_timestamps(data: bytes, out: Path) -> None:
    """python-pptx stamps zip entries with the current time; pin them so output is stable."""
    with (
        zipfile.ZipFile(io.BytesIO(data)) as src,
        zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as dst,
    ):
        for info in src.infolist():
            fixed = zipfile.ZipInfo(info.filename, date_time=ZIP_DATE_TIME)
            fixed.compress_type = zipfile.ZIP_DEFLATED
            fixed.external_attr = info.external_attr
            dst.writestr(fixed, src.read(info.filename))


# --- Entry point -------------------------------------------------------------------------


def main() -> None:
    png = make_png()
    PNG_PATH.write_bytes(png)
    make_pdf(png)
    make_pptx(png)
    for path in (PNG_PATH, PDF_PATH, PPTX_PATH):
        print(f"{path} ({path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
