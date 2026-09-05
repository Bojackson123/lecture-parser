"""P1-04 ``lecturenotes captions FILE``, P2-04 ``lecturenotes slides FILE``,
P3-04 ``lecturenotes render FILE [-o DIR]``, P4-04 ``lecturenotes align DECK
CAPTIONS``, P5-04 ``lecturenotes build PATHS...`` and P7-05 ``lecturenotes push
FILE --parent PAGE_ID``.

Debugging commands, not the product (plan §8: "bad notes are almost always bad
chunks"): ``captions`` prints one line per segment, ``[m:ss–m:ss] text``, or the
segments as JSON that ``Segment.model_validate`` accepts back; ``slides`` prints one
deck - titles, blocks in reading order, images found, notes on request - or the
``Deck`` as JSON that ``Deck.model_validate_json`` accepts back; ``render`` prints the
markdown one ``NoteWeek`` JSON renders to (or emits it to a directory with ``-o``), or
the ``RenderResult`` as JSON that ``RenderResult.model_validate_json`` accepts back;
``align`` prints the chunks one deck and one caption file align to, or the chunk list
as JSON that ``Chunk.model_validate`` accepts back. ``build`` is the product: pairing
ritual (§7.4), ``--dry-run`` chunking (§8), and the fake-driven real run that writes
the ``NoteWeek`` JSON ``render`` consumes (§7.1's tuning loop). Everything goes
through ``main([...])`` and ``capsys`` so the tests exercise exactly what the console
script runs; the build tests monkeypatch ``cli._make_client``, the one client seam,
and the push tests monkeypatch ``cli._make_transport``, the one transport seam, so no
test touches the network (plan §8).
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from textwrap import dedent

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

import lecturenotes
from lecturenotes import cli
from lecturenotes.align.boundaries import Chunk
from lecturenotes.cli import format_clock, main
from lecturenotes.emit.notion_api import FakeNotionTransport
from lecturenotes.generate.client import GenRequest, RecordedClient
from lecturenotes.ingest.slides import Deck, image_id
from lecturenotes.model import NoteLecture, NoteWeek, SourceRef
from lecturenotes.render.base import RenderResult

SEGMENT_COUNT = 22
SLIDE_COUNT = 3
# README decks table, transcribed - never read from the parser output.
PPTX_IMAGE_ID = "img-a63ae9b7dc5e9397"
PPTX_IMAGE_LINE = f"[image {PPTX_IMAGE_ID} 240x150 image/png]"


@pytest.fixture(scope="module")
def vtt_path(fixtures_dir: Path) -> str:
    return str(fixtures_dir / "captions" / "lecture01.vtt")


@pytest.fixture(scope="module")
def srt_path(fixtures_dir: Path) -> str:
    return str(fixtures_dir / "captions" / "lecture01.srt")


# --- existing behaviour is unchanged ----------------------------------------------


def test_version_flag_still_works(capsys: pytest.CaptureFixture[str]) -> None:
    assert main(["--version"]) == 0
    assert capsys.readouterr().out.strip() == f"lecturenotes {lecturenotes.__version__}"


def test_no_arguments_prints_help_and_returns_0(capsys: pytest.CaptureFixture[str]) -> None:
    assert main([]) == 0
    out = capsys.readouterr().out
    assert out.startswith("usage: lecturenotes")
    assert "captions" in out
    assert "slides" in out
    assert "render" in out
    assert "align" in out


# --- plain lines ------------------------------------------------------------------


def test_captions_prints_22_lines_first_and_last_anchored(
    vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["captions", vtt_path]) == 0
    captured = capsys.readouterr()
    lines = captured.out.splitlines()
    assert len(lines) == SEGMENT_COUNT
    assert lines[0].startswith("[0:01–0:26] welcome back")
    assert lines[-1].startswith("[8:40–9:05] that's it")
    assert captured.err == ""


def test_every_line_has_the_bracketed_span_prefix(
    vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["captions", vtt_path])
    for line in capsys.readouterr().out.splitlines():
        prefix, _, text = line.partition("] ")
        assert prefix.startswith("["), line
        start, sep, end = prefix[1:].partition("–")
        assert sep and start and end, line
        assert text and text == text.strip(), line


def test_srt_prints_the_same_lines_as_vtt(
    vtt_path: str, srt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["captions", vtt_path]) == 0
    from_vtt = capsys.readouterr().out
    assert main(["captions", srt_path]) == 0
    assert capsys.readouterr().out == from_vtt


@pytest.mark.parametrize(
    ("seconds", "expected"),
    [
        (0.0, "0:00"),
        (1.0, "0:01"),
        (26.0, "0:26"),
        (59.9, "0:59"),  # floored, never rounded up into the next second
        (60.0, "1:00"),
        (520.0, "8:40"),
        (599.0, "9:59"),
        (600.0, "10:00"),
        (3599.0, "59:59"),
        (3600.0, "1:00:00"),
        (3723.0, "1:02:03"),
        (36_000.0, "10:00:00"),
    ],
)
def test_format_clock(seconds: float, expected: str) -> None:
    assert format_clock(seconds) == expected


# --- --json -----------------------------------------------------------------------


def test_json_output_is_22_segment_dicts(srt_path: str, capsys: pytest.CaptureFixture[str]) -> None:
    assert main(["captions", "--json", srt_path]) == 0
    data = json.loads(capsys.readouterr().out)
    assert isinstance(data, list)
    assert len(data) == SEGMENT_COUNT
    assert all(set(d) == {"start_s", "end_s", "text"} for d in data)


def test_json_output_equals_the_hand_written_fixture(
    vtt_path: str, fixtures_dir: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    """The JSON uses the ``Segment`` field names, so it round-trips with the snapshot."""
    main(["captions", "--json", vtt_path])
    printed = json.loads(capsys.readouterr().out)
    expected_raw = (fixtures_dir / "captions" / "lecture01.segments.json").read_text("utf-8")
    assert printed == json.loads(expected_raw)


# --- merge knobs ------------------------------------------------------------------


def test_max_segment_s_is_forwarded(vtt_path: str, capsys: pytest.CaptureFixture[str]) -> None:
    """40 s is below the fixture's longest merged span (49 s), so the output must change."""
    assert main(["captions", vtt_path]) == 0
    default = capsys.readouterr().out
    assert main(["captions", "--max-segment-s", "40", vtt_path]) == 0
    assert capsys.readouterr().out != default
    assert main(["captions", "--max-segment-s", "60", vtt_path]) == 0
    assert capsys.readouterr().out == default


def test_max_gap_s_is_forwarded(tmp_path: Path, capsys: pytest.CaptureFixture[str]) -> None:
    """Two unterminated cues 10 s apart: the default 5 s gap flushes them as two
    segments, ``--max-gap-s 20`` joins them into one. (The lecture fixture cannot show
    this — every sentence that carries across cues there crosses a zero-second gap.)"""
    vtt = tmp_path / "gap.vtt"
    vtt.write_text(
        dedent(
            """
            WEBVTT

            00:00:00.000 --> 00:00:05.000
            first half

            00:00:15.000 --> 00:00:20.000
            second half.
            """
        ).strip(),
        encoding="utf-8",
    )
    assert main(["captions", str(vtt)]) == 0
    assert capsys.readouterr().out.splitlines() == [
        "[0:00–0:05] first half",
        "[0:15–0:20] second half.",
    ]
    assert main(["captions", "--max-gap-s", "20", str(vtt)]) == 0
    assert capsys.readouterr().out.splitlines() == ["[0:00–0:20] first half second half."]


# --- errors -----------------------------------------------------------------------


def test_unsupported_suffix_returns_2_with_stderr_only(
    fixtures_dir: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    pdf = str(fixtures_dir / "decks" / "lecture01.pdf")
    assert main(["captions", pdf]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert ".pdf" in captured.err
    assert "Traceback" not in captured.err


def test_missing_file_returns_2_with_stderr_only(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    missing = str(tmp_path / "nope.vtt")
    assert main(["captions", missing]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "nope.vtt" in captured.err
    assert "Traceback" not in captured.err


def test_malformed_captions_return_2_with_the_line_number(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    bad = tmp_path / "bad.vtt"
    bad.write_text("WEBVTT\n\nthis is not a timing line\nhello\n", encoding="utf-8")
    assert main(["captions", str(bad)]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "line 3" in captured.err
    assert "Traceback" not in captured.err


# =====================================================================================
# P2-04: ``lecturenotes slides FILE``
# =====================================================================================


@pytest.fixture(scope="module")
def pptx_path(fixtures_dir: Path) -> str:
    return str(fixtures_dir / "decks" / "lecture01.pptx")


@pytest.fixture(scope="module")
def pdf_path(fixtures_dir: Path) -> str:
    return str(fixtures_dir / "decks" / "lecture01.pdf")


def _slide_headers(out: str) -> list[str]:
    return [line for line in out.splitlines() if line.startswith("--- slide")]


def _image_lines(out: str) -> list[str]:
    return [line for line in out.splitlines() if line.startswith("[image ")]


# --- plain lines ------------------------------------------------------------------


def test_slides_prints_3_slide_headers_with_titles(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["slides", pptx_path]) == 0
    captured = capsys.readouterr()
    assert _slide_headers(captured.out) == [
        "--- slide 1: Markov Decision Processes",
        "--- slide 2: The Bellman Equation",
        "--- slide 3: Value Iteration",
    ]
    assert captured.err == ""


def test_slides_prints_the_left_column_before_the_right(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    """Slide 2's reading order at a glance: the last left-column row precedes the
    right column's heading."""
    main(["slides", pptx_path])
    lines = capsys.readouterr().out.splitlines()
    assert lines.index("gamma: discount factor") < lines.index("Intuition")


def test_slides_separates_blocks_with_one_blank_line(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    """Slide 2 has two blocks: six left-column lines, a blank, six right-column lines."""
    main(["slides", pptx_path])
    lines = capsys.readouterr().out.splitlines()
    start = lines.index("--- slide 2: The Bellman Equation")
    end = lines.index("--- slide 3: Value Iteration")
    body = [line for line in lines[start + 1 : end] if not line.startswith("[")]
    while body and body[-1] == "":  # the slide separator, if the command prints one
        body.pop()
    assert body[0] == "Equation"
    assert body[6] == ""
    assert body[7] == "Intuition"
    assert body.count("") == 1
    assert len(body) == 13


def test_slides_lists_the_figure_as_an_image_line_after_slide_3s_text(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["slides", pptx_path])
    out = capsys.readouterr().out
    assert _image_lines(out) == [PPTX_IMAGE_LINE]
    lines = out.splitlines()
    last_step = "4. Read off the greedy policy pi(s) = argmax_a [ ... ]"
    assert lines.index("--- slide 3: Value Iteration") < lines.index(last_step)
    assert lines.index(last_step) < lines.index(PPTX_IMAGE_LINE)


def test_slides_prints_no_notes_by_default(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["slides", pptx_path])
    assert "[notes]" not in capsys.readouterr().out


def test_slides_notes_flag_adds_one_notes_line_per_slide(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["slides", "--notes", pptx_path]) == 0
    lines = capsys.readouterr().out.splitlines()
    notes = [line for line in lines if line.startswith("[notes] ")]
    assert len(notes) == SLIDE_COUNT
    assert "this will be on the exam" in notes[1]
    # Each notes line sits inside its own slide, after the text and images.
    headers = _slide_headers("\n".join(lines))
    for header, note in zip(headers, notes, strict=True):
        assert lines.index(header) < lines.index(note)
    assert lines.index(PPTX_IMAGE_LINE) < lines.index(notes[2])
    assert lines.index(notes[0]) < lines.index(headers[1])
    assert lines.index(notes[1]) < lines.index(headers[2])


def test_slides_pdf_prints_the_same_headers_and_never_notes(
    pptx_path: str, pdf_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["slides", pptx_path]) == 0
    from_pptx = _slide_headers(capsys.readouterr().out)
    assert main(["slides", "--notes", pdf_path]) == 0
    out = capsys.readouterr().out
    assert _slide_headers(out) == from_pptx
    assert "[notes]" not in out
    # The PDF figure is re-encoded (a different id) but is the same 240x150 PNG.
    image_lines = _image_lines(out)
    assert len(image_lines) == 1
    assert image_lines[0].startswith("[image img-")
    assert image_lines[0].endswith(" 240x150 image/png]")
    assert image_lines[0] != PPTX_IMAGE_LINE


def test_slides_pdf_and_pptx_print_the_same_text_lines(
    pptx_path: str, pdf_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    """The cross-format assertion, seen from the shell: everything but the image line."""

    def text_only(out: str) -> list[str]:
        return [line for line in out.splitlines() if not line.startswith("[image ")]

    assert main(["slides", pptx_path]) == 0
    from_pptx = text_only(capsys.readouterr().out)
    assert main(["slides", pdf_path]) == 0
    assert text_only(capsys.readouterr().out) == from_pptx


# --- --json -----------------------------------------------------------------------


def test_slides_json_validates_back_to_a_3_slide_deck(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["slides", "--json", pptx_path]) == 0
    out = capsys.readouterr().out
    deck = Deck.model_validate_json(out)
    assert len(deck.slides) == SLIDE_COUNT
    assert [a.id for a in deck.assets] == [PPTX_IMAGE_ID]
    assert json.loads(out)["slides"][2]["notes"]  # --json is the complete view


def test_slides_json_equals_the_hand_written_fixture(
    pptx_path: str, fixtures_dir: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["slides", "--json", pptx_path])
    printed = Deck.model_validate_json(capsys.readouterr().out)
    expected_raw = (fixtures_dir / "decks" / "lecture01.deck.json").read_text("utf-8")
    expected = Deck.model_validate_json(expected_raw)
    # source is the path as given (absolute here, repo-relative in the fixture).
    assert printed.source.endswith(expected.source)
    assert printed.model_copy(update={"source": expected.source}) == expected


# --- ad-hoc decks: [hidden], untitled, --min-px, [recurring] ----------------------------


def _png(tmp_path: Path, name: str, size: tuple[int, int], colour: str) -> Path:
    path = tmp_path / f"{name}.png"
    PILImage.new("RGB", size, colour).save(path)
    return path


def _pptx(tmp_path: Path, per_slide: list[list[Path]], hidden: tuple[int, ...] = ()) -> str:
    """One blank, untitled slide per entry with its pictures side by side."""
    blank = 6
    prs = Presentation()
    for number, pictures in enumerate(per_slide, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[blank])
        for i, picture in enumerate(pictures):
            slide.shapes.add_picture(str(picture), Inches(1 + 3 * i), Inches(1))
        if number in hidden:
            slide._element.set("show", "0")
    path = tmp_path / "adhoc.pptx"
    prs.save(str(path))
    return str(path)


def test_slides_untitled_and_hidden_headers(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    deck = _pptx(tmp_path, [[], [], []], hidden=(2,))
    assert main(["slides", deck]) == 0
    assert _slide_headers(capsys.readouterr().out) == [
        "--- slide 1",
        "--- slide 2 [hidden]",
        "--- slide 3",
    ]


def test_slides_min_px_is_forwarded(tmp_path: Path, capsys: pytest.CaptureFixture[str]) -> None:
    tiny = _png(tmp_path, "tiny", (8, 8), "red")
    deck = _pptx(tmp_path, [[tiny]])
    assert main(["slides", deck]) == 0
    assert _image_lines(capsys.readouterr().out) == []
    assert main(["slides", "--min-px", "4", deck]) == 0
    image_lines = _image_lines(capsys.readouterr().out)
    assert len(image_lines) == 1
    assert image_lines[0].endswith(" 8x8 image/png]")


def test_slides_recurring_logo_is_listed_once_after_the_last_slide(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    logo = _png(tmp_path, "logo", (64, 48), "navy")
    figure = _png(tmp_path, "figure", (64, 48), "olive")
    deck = _pptx(tmp_path, [[logo], [logo, figure], [logo], []])
    assert main(["slides", deck]) == 0
    out = capsys.readouterr().out
    lines = out.splitlines()
    recurring = [line for line in lines if line.startswith("[recurring] ")]
    assert len(recurring) == 1
    logo_id = image_id(logo.read_bytes())
    assert recurring[0].split() == ["[recurring]", logo_id, "64x48", "image/png"]
    assert lines.index("--- slide 4") < lines.index(recurring[0])
    # The logo is nowhere in an [image ...] line; the figure on slide 2 still is.
    image_lines = _image_lines(out)
    assert len(image_lines) == 1
    assert image_id(figure.read_bytes()) in image_lines[0]
    assert lines.index("--- slide 2") < lines.index(image_lines[0]) < lines.index("--- slide 3")


def test_slides_prints_no_recurring_line_when_nothing_recurs(
    pptx_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["slides", pptx_path])
    assert "[recurring]" not in capsys.readouterr().out


# --- errors -----------------------------------------------------------------------


def test_slides_unsupported_suffix_returns_2_with_stderr_only(
    vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["slides", vtt_path]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes slides: ")
    assert ".vtt" in captured.err
    assert "Traceback" not in captured.err


def test_slides_missing_file_returns_2_with_stderr_only(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    missing = str(tmp_path / "nope.pptx")
    assert main(["slides", missing]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "nope.pptx" in captured.err
    assert "Traceback" not in captured.err


def test_slides_garbage_deck_returns_2_naming_the_file(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    bad = tmp_path / "bad.pdf"
    bad.write_bytes(b"not a pdf\n")
    assert main(["slides", str(bad)]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "bad.pdf" in captured.err
    assert "Traceback" not in captured.err


# =====================================================================================
# P3-04: ``lecturenotes render FILE [-o DIR]``
# =====================================================================================


@pytest.fixture(scope="module")
def week_json_path(fixtures_dir: Path) -> str:
    return str(fixtures_dir / "notes" / "week01.json")


@pytest.fixture(scope="module")
def expected_markdown(fixtures_dir: Path) -> str:
    """Hand-written (P3-02) — bytes, not read_text, so newline handling can't lie."""
    return (fixtures_dir / "notes" / "week01.md").read_bytes().decode("utf-8")


# --- stdout mode ------------------------------------------------------------------


def test_render_prints_one_document_under_its_name(
    week_json_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", week_json_path]) == 0
    captured = capsys.readouterr()
    headers = [line for line in captured.out.splitlines() if line.startswith("--- ")]
    assert headers == ["--- cs-rl-101-w01.md"]
    assert "$$" in captured.out
    assert "> **EXAM**" in captured.out
    assert "[2:31–4:28]" in captured.out  # the slide-less anchor
    assert captured.err == ""


def test_render_text_after_the_header_equals_the_expected_markdown(
    week_json_path: str, expected_markdown: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["render", week_json_path])
    header, _, body = capsys.readouterr().out.partition("\n")
    assert header == "--- cs-rl-101-w01.md"
    assert body == expected_markdown


# --- -o DIR -----------------------------------------------------------------------


def test_render_out_writes_page_and_asset_and_prints_nothing(
    week_json_path: str,
    expected_markdown: str,
    fixtures_dir: Path,
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    assert main(["render", week_json_path, "-o", str(tmp_path)]) == 0
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == ""
    page = tmp_path / "cs-rl-101-w01.md"
    assert page.read_bytes().decode("utf-8") == expected_markdown
    copied = tmp_path / "assets" / "fig-value-iteration-convergence.png"
    original = fixtures_dir / "decks" / "value_iteration.png"
    assert copied.read_bytes() == original.read_bytes()


# --- --json -----------------------------------------------------------------------


def test_render_json_revalidates_to_one_document_and_one_asset(
    week_json_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", "--json", week_json_path]) == 0
    result = RenderResult.model_validate_json(capsys.readouterr().out)
    assert [d.name for d in result.documents] == ["cs-rl-101-w01.md"]
    assert [a.id for a in result.assets] == ["fig-value-iteration-convergence"]


# --- --format (P6-03) -------------------------------------------------------------


@pytest.fixture(scope="module")
def expected_anki_deck(fixtures_dir: Path) -> str:
    """Hand-written (P6-01) — bytes, not read_text, so newline handling can't lie."""
    return (fixtures_dir / "notes" / "week01.anki.txt").read_bytes().decode("utf-8")


def test_render_format_anki_prints_the_expected_deck(
    week_json_path: str, expected_anki_deck: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", week_json_path, "--format", "anki"]) == 0
    captured = capsys.readouterr()
    header, _, body = captured.out.partition("\n")
    assert header == "--- cs-rl-101-w01.txt"
    assert body == expected_anki_deck
    assert len([line for line in body.splitlines() if line.startswith("#")]) == 6
    assert len([line for line in body.splitlines() if "\t" in line]) == 8
    assert captured.err == ""


def test_render_format_markdown_prints_exactly_what_no_flag_prints(
    week_json_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", week_json_path]) == 0
    default_out = capsys.readouterr().out
    assert main(["render", week_json_path, "--format", "markdown"]) == 0
    assert capsys.readouterr().out == default_out


def test_render_format_anki_out_writes_deck_and_no_assets_dir(
    week_json_path: str,
    expected_anki_deck: str,
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    """The deck's manifest is empty, and the emitter creates ``assets/`` only when the
    manifest is non-empty (P3-03) — so an Anki emit is exactly one file."""
    assert main(["render", week_json_path, "--format", "anki", "-o", str(tmp_path)]) == 0
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == ""
    deck = tmp_path / "cs-rl-101-w01.txt"
    assert deck.read_bytes().decode("utf-8") == expected_anki_deck
    assert not (tmp_path / "assets").exists()


def test_render_format_anki_json_revalidates_to_one_document_and_no_assets(
    week_json_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    """--json stays format-agnostic: it dumps ``RenderResult`` whatever produced it."""
    assert main(["render", "--json", week_json_path, "--format", "anki"]) == 0
    result = RenderResult.model_validate_json(capsys.readouterr().out)
    assert [d.name for d in result.documents] == ["cs-rl-101-w01.txt"]
    assert list(result.assets) == []


# --- errors -----------------------------------------------------------------------


def test_render_non_json_file_returns_2_with_stderr_only(
    vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", vtt_path]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes render: ")
    assert "Traceback" not in captured.err


def test_render_wrong_shape_json_returns_2_with_stderr_only(
    fixtures_dir: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    """Valid JSON that is a ``Deck``, not a ``NoteWeek`` — pydantic's ValidationError
    is a ValueError, so it gets the uniform error line, not a traceback."""
    deck_json = str(fixtures_dir / "decks" / "lecture01.deck.json")
    assert main(["render", deck_json]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes render: ")
    assert "Traceback" not in captured.err


def test_render_missing_file_returns_2_with_stderr_only(
    tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    missing = str(tmp_path / "nope.json")
    assert main(["render", missing]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "nope.json" in captured.err
    assert "Traceback" not in captured.err


# =====================================================================================
# P4-04: ``lecturenotes align DECK CAPTIONS``
# =====================================================================================

CHUNK_COUNT = 4
# The slide → time map in tests/fixtures/README.md, transcribed — never read from the
# aligner's output. The gap header says "(no slide)", a fact about the material, not
# "gap", which is pipeline jargon that reads as an error.
EXPECTED_CHUNK_HEADERS = [
    "--- slide 1: Markov Decision Processes [0:01–2:29]",
    "--- (no slide) [2:31–4:28]",
    "--- slide 2: The Bellman Equation [4:31–6:59]",
    "--- slide 3: Value Iteration [7:01–9:05]",
]


def _chunk_headers(out: str) -> list[str]:
    return [line for line in out.splitlines() if line.startswith("--- ")]


def _segment_lines(out: str) -> list[str]:
    return [line for line in out.splitlines() if line.startswith("  [")]


# --- plain lines ------------------------------------------------------------------


def test_align_prints_4_chunk_headers_in_order(
    pdf_path: str, vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["align", pdf_path, vtt_path]) == 0
    captured = capsys.readouterr()
    assert _chunk_headers(captured.out) == EXPECTED_CHUNK_HEADERS
    assert captured.err == ""


def test_align_indents_all_22_segments_under_their_chunks(
    pdf_path: str, vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    main(["align", pdf_path, vtt_path])
    lines = _segment_lines(capsys.readouterr().out)
    assert len(lines) == SEGMENT_COUNT
    assert lines[0].startswith("  [0:01–0:26] welcome back")
    assert lines[-1].startswith("  [8:40–9:05] that's it")


def test_align_pptx_and_srt_print_identical_stdout(
    pptx_path: str, pdf_path: str, vtt_path: str, srt_path: str,
    capsys: pytest.CaptureFixture[str],
) -> None:
    """The cross-format invariant, end to end from the shell."""
    assert main(["align", pdf_path, vtt_path]) == 0
    from_pdf_vtt = capsys.readouterr().out
    assert main(["align", pptx_path, srt_path]) == 0
    assert capsys.readouterr().out == from_pdf_vtt


def test_align_min_gap_s_is_forwarded(
    pdf_path: str, vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    """1000 s exceeds the dice detour's span, so no stretch qualifies as a gap."""
    assert main(["align", "--min-gap-s", "1000", pdf_path, vtt_path]) == 0
    headers = _chunk_headers(capsys.readouterr().out)
    assert len(headers) == 3
    assert not any("(no slide)" in header for header in headers)


# --- --json -----------------------------------------------------------------------


def test_align_json_revalidates_to_4_chunks_with_a_gap(
    pdf_path: str, vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["align", "--json", pdf_path, vtt_path]) == 0
    data = json.loads(capsys.readouterr().out)
    assert isinstance(data, list)
    assert len(data) == CHUNK_COUNT
    chunks = [Chunk.model_validate(element) for element in data]
    assert data[1]["slides"] is None
    assert sum(len(chunk.segments) for chunk in chunks) == SEGMENT_COUNT


# --- errors -----------------------------------------------------------------------


def test_align_wrong_deck_returns_2_with_stderr_only(
    week_json_path: str, vtt_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["align", week_json_path, vtt_path]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes align: ")
    assert "Traceback" not in captured.err


def test_align_missing_captions_returns_2_with_stderr_only(
    pdf_path: str, tmp_path: Path, capsys: pytest.CaptureFixture[str]
) -> None:
    missing = str(tmp_path / "nope.vtt")
    assert main(["align", pdf_path, missing]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "nope.vtt" in captured.err
    assert "Traceback" not in captured.err


# =====================================================================================
# P5-04: ``lecturenotes build PATHS... --course TEXT --week N``
# =====================================================================================

WEEK_ID = "cs-rl-101-w01"
MEDIA_NAME = f"{PPTX_IMAGE_ID}.png"
COURSE_ARGS = ["--course", "CS-RL-101", "--week", "1"]


@pytest.fixture(scope="module")
def responses_path(fixtures_dir: Path) -> Path:
    return fixtures_dir / "generate" / "lecture01.responses.json"


@pytest.fixture(scope="module")
def notes_fixture_path(fixtures_dir: Path) -> Path:
    return fixtures_dir / "generate" / "lecture01.notes.json"


@pytest.fixture
def no_client(monkeypatch: pytest.MonkeyPatch) -> None:
    """Any client construction is a test failure — and no key must ever be consulted."""

    def boom(model: str) -> RecordedClient:
        raise AssertionError("a client was constructed")

    monkeypatch.setattr(cli, "_make_client", boom)
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)


class _CountingRecorded:
    """The recorded fake behind a shared call counter — the CLI-level §7.1 budget."""

    def __init__(self, path: Path, counter: list[int]) -> None:
        self._inner = RecordedClient(path)
        self.model = self._inner.model
        self._counter = counter

    def complete(self, request: GenRequest) -> str:
        self._counter[0] += 1
        return self._inner.complete(request)


@pytest.fixture
def complete_calls(
    responses_path: Path, monkeypatch: pytest.MonkeyPatch
) -> list[int]:
    """Patch the client seam with the counting recorded fake; yield the call counter."""
    counter = [0]
    monkeypatch.setattr(
        cli, "_make_client", lambda model: _CountingRecorded(responses_path, counter)
    )
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    return counter


class _FakeTty:
    def isatty(self) -> bool:
        return True


def _pairing_lines(out: str) -> list[str]:
    return [line for line in out.splitlines() if line.lstrip().startswith("lec")]


# --- --dry-run ---------------------------------------------------------------------


def test_build_dry_run_prints_pairing_then_4_chunks_with_no_client(
    pptx_path: str, vtt_path: str, no_client: None, capsys: pytest.CaptureFixture[str]
) -> None:
    """The plan §6 dry-run criterion: pairing + chunking, before any client exists."""
    assert main(["build", pptx_path, vtt_path, *COURSE_ARGS, "--dry-run"]) == 0
    captured = capsys.readouterr()
    assert captured.err == ""
    pairing = _pairing_lines(captured.out)
    assert len(pairing) == 1
    assert "lec01" in pairing[0]
    assert "lecture01.pptx" in pairing[0]
    assert "lecture01.vtt" in pairing[0]
    assert _chunk_headers(captured.out) == EXPECTED_CHUNK_HEADERS
    # The pairing comes first; the chunks follow in the align command's format.
    lines = captured.out.splitlines()
    assert lines.index(pairing[0]) < lines.index(EXPECTED_CHUNK_HEADERS[0])


def test_build_dry_run_forwards_min_words(
    pptx_path: str, vtt_path: str, no_client: None, capsys: pytest.CaptureFixture[str]
) -> None:
    """At a 200-word floor the three slide chunks merge into one; the gap fences."""
    argv = ["build", pptx_path, vtt_path, *COURSE_ARGS, "--dry-run", "--min-words", "200"]
    assert main(argv) == 0
    headers = _chunk_headers(capsys.readouterr().out)
    assert len(headers) < len(EXPECTED_CHUNK_HEADERS)
    assert sum("(no slide)" in header for header in headers) == 1


# --- pairing -----------------------------------------------------------------------


def _sources_dir(tmp_path: Path, pptx_path: str, vtt_path: str, lectures: int) -> Path:
    """A directory of ``lectureNN.pptx``/``lectureNN.vtt`` copies of the fixtures."""
    sources = tmp_path / "sources"
    sources.mkdir()
    for n in range(1, lectures + 1):
        shutil.copy(pptx_path, sources / f"lecture{n:02d}.pptx")
        shutil.copy(vtt_path, sources / f"lecture{n:02d}.vtt")
    return sources


def test_build_pairs_a_directory_in_sorted_filename_order(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    no_client: None,
    capsys: pytest.CaptureFixture[str],
) -> None:
    sources = _sources_dir(tmp_path, pptx_path, vtt_path, lectures=2)
    assert main(["build", str(sources), *COURSE_ARGS, "--dry-run"]) == 0
    pairing = _pairing_lines(capsys.readouterr().out)
    assert len(pairing) == 2
    assert "lec01" in pairing[0] and "lecture01.pptx" in pairing[0]
    assert "lecture01.vtt" in pairing[0]
    assert "lec02" in pairing[1] and "lecture02.pptx" in pairing[1]
    assert "lecture02.vtt" in pairing[1]


def test_build_unequal_counts_exit_2_listing_both_sides(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    no_client: None,
    capsys: pytest.CaptureFixture[str],
) -> None:
    sources = _sources_dir(tmp_path, pptx_path, vtt_path, lectures=2)
    (sources / "lecture02.vtt").unlink()
    assert main(["build", str(sources), *COURSE_ARGS, "--dry-run"]) == 2
    captured = capsys.readouterr()
    assert captured.err.startswith("lecturenotes build: ")
    assert "lecture01.pptx" in captured.err and "lecture02.pptx" in captured.err
    assert "lecture01.vtt" in captured.err
    assert "Traceback" not in captured.err


def test_build_unknown_suffix_exits_2(
    vtt_path: str, tmp_path: Path, no_client: None, capsys: pytest.CaptureFixture[str]
) -> None:
    stray = tmp_path / "syllabus.txt"
    stray.write_text("not lecture material\n", encoding="utf-8")
    assert main(["build", str(stray), vtt_path, *COURSE_ARGS, "--dry-run"]) == 2
    captured = capsys.readouterr()
    assert "syllabus.txt" in captured.err
    assert "Traceback" not in captured.err


# --- confirmation ------------------------------------------------------------------


def test_build_decline_exits_1_after_printing_the_pairing(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    no_client: None,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    monkeypatch.setattr("sys.stdin", _FakeTty())
    monkeypatch.setattr("builtins.input", lambda prompt="": "n")
    argv = ["build", pptx_path, vtt_path, *COURSE_ARGS, "-o", str(tmp_path / "out")]
    assert main(argv) == 1  # no_client asserts no client was ever constructed
    out = capsys.readouterr().out
    assert len(_pairing_lines(out)) == 1
    assert not (tmp_path / "out").exists()


def test_build_confirm_y_proceeds(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    complete_calls: list[int],
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr("sys.stdin", _FakeTty())
    monkeypatch.setattr("builtins.input", lambda prompt="": "y")
    argv = ["build", pptx_path, vtt_path, *COURSE_ARGS, "-o", str(tmp_path / "out")]
    assert main(argv) == 0
    assert (tmp_path / "out" / f"{WEEK_ID}.json").exists()


def test_build_non_tty_without_yes_exits_2_naming_the_flag(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    no_client: None,
    capsys: pytest.CaptureFixture[str],
) -> None:
    """Scripts must never hang on a hidden prompt (pytest's stdin is not a TTY)."""
    argv = ["build", pptx_path, vtt_path, *COURSE_ARGS, "-o", str(tmp_path / "out")]
    assert main(argv) == 2
    captured = capsys.readouterr()
    assert "--yes" in captured.err
    assert "Traceback" not in captured.err
    assert not (tmp_path / "out").exists()


# --- the real run (recorded fake) --------------------------------------------------


def _build(pptx_path: str, vtt_path: str, out: Path) -> int:
    return main(["build", pptx_path, vtt_path, *COURSE_ARGS, "--yes", "-o", str(out)])


def test_build_real_run_writes_the_week_json_and_media(
    pptx_path: str,
    vtt_path: str,
    notes_fixture_path: Path,
    tmp_path: Path,
    complete_calls: list[int],
) -> None:
    """The plan §6 criterion, fake-driven: one valid ``NoteWeek``, the fixture lecture
    with the command's own paths as its source, and the referenced media file."""
    assert _build(pptx_path, vtt_path, tmp_path) == 0
    raw = (tmp_path / f"{WEEK_ID}.json").read_bytes()
    assert b"\r" not in raw and raw.endswith(b"\n")  # UTF-8, LF, newline-terminated
    week = NoteWeek.model_validate_json(raw.decode("utf-8"))
    assert week.id == WEEK_ID
    assert week.course == "CS-RL-101"
    assert week.week_number == 1
    expected = NoteLecture.model_validate_json(
        notes_fixture_path.read_text(encoding="utf-8")
    ).model_copy(
        update={
            "source": SourceRef(
                deck_path=Path(pptx_path).as_posix(),
                caption_path=Path(vtt_path).as_posix(),
            )
        }
    )
    assert week.lectures == [expected]
    assert (tmp_path / "media" / MEDIA_NAME).exists()


def test_build_output_renders(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    complete_calls: list[int],
    capsys: pytest.CaptureFixture[str],
) -> None:
    """§7.1's tuning loop end to end: ``build`` then ``render``, two commands."""
    assert _build(pptx_path, vtt_path, tmp_path) == 0
    capsys.readouterr()
    assert main(["render", str(tmp_path / f"{WEEK_ID}.json")]) == 0
    captured = capsys.readouterr()
    headers = [line for line in captured.out.splitlines() if line.startswith("--- ")]
    assert headers == [f"--- {WEEK_ID}.md"]
    assert captured.err == ""


def test_build_second_run_is_served_entirely_from_the_cache(
    pptx_path: str,
    vtt_path: str,
    tmp_path: Path,
    complete_calls: list[int],
) -> None:
    """The default ``<out>/.cache`` persisted all five responses: zero new
    ``complete`` calls, byte-identical output files (§7.1 through the CLI)."""
    assert _build(pptx_path, vtt_path, tmp_path) == 0
    assert complete_calls[0] == 5
    first_json = (tmp_path / f"{WEEK_ID}.json").read_bytes()
    first_media = (tmp_path / "media" / MEDIA_NAME).read_bytes()
    assert _build(pptx_path, vtt_path, tmp_path) == 0
    assert complete_calls[0] == 5
    assert (tmp_path / f"{WEEK_ID}.json").read_bytes() == first_json
    assert (tmp_path / "media" / MEDIA_NAME).read_bytes() == first_media


# =====================================================================================
# P7-05: ``render --format notion`` and ``lecturenotes push FILE --parent PAGE_ID``
# =====================================================================================

NOTION_DOC = f"{WEEK_ID}.notion.json"
NOTION_TITLE = "CS-RL-101 — Week 1"
FIGURE_PNG = "fig-value-iteration-convergence.png"


@pytest.fixture(scope="module")
def expected_notion_payload(fixtures_dir: Path) -> str:
    """Hand-written (P7-01) — bytes, not read_text, so newline handling can't lie."""
    return (fixtures_dir / "notes" / "week01.notion.json").read_bytes().decode("utf-8")


# --- render --format notion ---------------------------------------------------------


def test_render_format_notion_prints_the_expected_payload(
    week_json_path: str, expected_notion_payload: str, capsys: pytest.CaptureFixture[str]
) -> None:
    assert main(["render", week_json_path, "--format", "notion"]) == 0
    captured = capsys.readouterr()
    header, _, body = captured.out.partition("\n")
    assert header == f"--- {NOTION_DOC}"
    assert body == expected_notion_payload
    assert set(json.loads(body)) == {"page", "payloads"}
    assert captured.err == ""


def test_render_format_notion_out_writes_the_document_and_the_figure(
    week_json_path: str,
    expected_notion_payload: str,
    fixtures_dir: Path,
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    """Unlike anki's empty manifest, notion's ``-o`` copies the figure too — the
    offline view of exactly what ``push`` would upload."""
    assert main(["render", week_json_path, "--format", "notion", "-o", str(tmp_path)]) == 0
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == ""
    document = tmp_path / NOTION_DOC
    assert document.read_bytes().decode("utf-8") == expected_notion_payload
    copied = tmp_path / "assets" / FIGURE_PNG
    original = fixtures_dir / "decks" / "value_iteration.png"
    assert copied.read_bytes() == original.read_bytes()


# --- push wiring and errors ---------------------------------------------------------


@pytest.fixture
def no_transport(monkeypatch: pytest.MonkeyPatch) -> None:
    """Any transport construction is a test failure (the P5-01 doctrine)."""

    def boom(token: str) -> FakeNotionTransport:
        raise AssertionError("a transport was constructed")

    monkeypatch.setattr(cli, "_make_transport", boom)
    monkeypatch.delenv("NOTION_TOKEN", raising=False)


class _TransportSeam:
    """The factory the tests inject: one shared fake, tokens recorded."""

    def __init__(self) -> None:
        self.fake = FakeNotionTransport()
        self.tokens: list[str] = []

    def __call__(self, token: str) -> FakeNotionTransport:
        self.tokens.append(token)
        return self.fake


@pytest.fixture
def transport_seam(monkeypatch: pytest.MonkeyPatch) -> _TransportSeam:
    seam = _TransportSeam()
    monkeypatch.setattr(cli, "_make_transport", seam)
    monkeypatch.setenv("NOTION_TOKEN", "secret-token")
    return seam


def test_push_without_parent_is_an_argparse_error(
    week_json_path: str, capsys: pytest.CaptureFixture[str]
) -> None:
    with pytest.raises(SystemExit) as excinfo:
        main(["push", week_json_path])
    assert excinfo.value.code == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert "--parent" in captured.err


def test_push_missing_file_returns_2_with_stderr_only(
    tmp_path: Path, no_transport: None, capsys: pytest.CaptureFixture[str]
) -> None:
    missing = str(tmp_path / "nope.json")
    assert main(["push", missing, "--parent", "x"]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes push: ")
    assert "nope.json" in captured.err
    assert "Traceback" not in captured.err


def test_push_wrong_shape_json_returns_2_with_stderr_only(
    fixtures_dir: Path, no_transport: None, capsys: pytest.CaptureFixture[str]
) -> None:
    """The standard render-path error contract: a ``Deck``, not a ``NoteWeek``."""
    deck_json = str(fixtures_dir / "decks" / "lecture01.deck.json")
    assert main(["push", deck_json, "--parent", "x"]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err.startswith("lecturenotes push: ")
    assert "Traceback" not in captured.err


def test_push_without_token_exits_2_and_constructs_no_transport(
    week_json_path: str, no_transport: None, capsys: pytest.CaptureFixture[str]
) -> None:
    """``NOTION_TOKEN`` is read at run time, after validation — and its absence must
    fail cleanly before any transport exists (the P5-01 doctrine, asserted the same
    way as ``no_client``: the boom factory would raise)."""
    assert main(["push", week_json_path, "--parent", "x"]) == 2
    captured = capsys.readouterr()
    assert captured.out == ""
    assert captured.err == "lecturenotes push: NOTION_TOKEN is not set\n"


# --- push against the injected fake -------------------------------------------------


def _push_fixture_week(week_json_path: str, repo_root: Path) -> int:
    """The fixture's asset sources are repo-root-relative (the P3-04 quirk), so the
    fixture push overrides the week-JSON-directory default with ``--asset-root``."""
    return main(
        ["push", week_json_path, "--parent", "parent-1", "--asset-root", str(repo_root)]
    )


def test_push_fresh_emit_runs_the_full_sequence_and_prints_a_summary(
    week_json_path: str,
    repo_root: Path,
    transport_seam: _TransportSeam,
    capsys: pytest.CaptureFixture[str],
) -> None:
    assert _push_fixture_week(week_json_path, repo_root) == 0
    fake = transport_seam.fake
    assert [call[0] for call in fake.calls] == [
        "find_child_page",
        "create_page",
        "upload_file",
        "append_children",
    ]
    assert fake.calls[0][1:] == ("parent-1", NOTION_TITLE)
    (page_id, appended), = fake.appended
    assert page_id == "page-1"
    assert appended[0]["type"] == "heading_1"
    name, media_type, data = fake.uploaded["upload-1"]
    assert name == FIGURE_PNG
    assert media_type == "image/png"
    assert data == (repo_root / "tests/fixtures/decks/value_iteration.png").read_bytes()
    assert transport_seam.tokens == ["secret-token"]
    captured = capsys.readouterr()
    assert captured.err == ""
    assert captured.out == f'pushed "{NOTION_TITLE}": 1 payload(s), 1 asset(s)\n'


def test_push_twice_reemits_to_the_same_page_with_no_second_create(
    week_json_path: str,
    repo_root: Path,
    transport_seam: _TransportSeam,
    capsys: pytest.CaptureFixture[str],
) -> None:
    """§7.2 at the CLI: the second push finds the page by title and never creates."""
    assert _push_fixture_week(week_json_path, repo_root) == 0
    first_calls = len(transport_seam.fake.calls)
    assert _push_fixture_week(week_json_path, repo_root) == 0
    second = [call[0] for call in transport_seam.fake.calls[first_calls:]]
    assert second == ["find_child_page", "list_children", "upload_file", "append_children"]
    assert [page_id for page_id, _ in transport_seam.fake.appended] == ["page-1", "page-1"]


def test_push_asset_root_defaults_to_the_week_jsons_directory(
    week_json_path: str, tmp_path: Path, transport_seam: _TransportSeam
) -> None:
    """The P5-03 layout observed end-to-end: a week JSON with its ``media/`` beside it
    pushes those bytes with no flag, wherever the process runs from."""
    week = json.loads(Path(week_json_path).read_text(encoding="utf-8"))
    week["lectures"][0]["assets"][0]["source"] = "media/value_iteration.png"
    moved = tmp_path / f"{WEEK_ID}.json"
    moved.write_text(json.dumps(week), encoding="utf-8")
    media = tmp_path / "media"
    media.mkdir()
    sentinel = b"bytes-from-the-week-dirs-media"
    (media / "value_iteration.png").write_bytes(sentinel)

    assert main(["push", str(moved), "--parent", "parent-1"]) == 0
    ((_upload_id, (name, _media_type, data)),) = transport_seam.fake.uploaded.items()
    assert name == FIGURE_PNG
    assert data == sentinel
