"""P2-03: slide images - PDF extraction, the size filter, the recurring-image rule, groups.

The fixture deck has one figure on slide 3, so the two format-agnostic rules
(``min_px``, recurring) are exercised on ad-hoc PPTX decks built with python-pptx under
``tmp_path`` with pictures generated in-test by Pillow - the slide-side analogue of
Phase 1's inline silence-gap cues. The canonical deck is never mutated (ticket decision).
Broken PDF images are fixture page 3 copied with ``pypdf.PdfWriter`` and its image
stream corrupted two ways.
"""

from __future__ import annotations

from collections.abc import Callable
from pathlib import Path

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from pypdf.generic import DictionaryObject, NameObject

from lecturenotes.ingest.slides import (
    Deck,
    Slide,
    SlideImage,
    TextBlock,
    _media_type_ok,
    drop_small_images,
    image_id,
    ingest_slides,
    parse_pdf,
    parse_pptx,
    set_aside_recurring_images,
)

# README decks table, transcribed - never read from the parser output.
PPTX_IMAGE_ID = "img-a63ae9b7dc5e9397"
FIGURE_SIZE = (240, 150)

BLANK = 6
PICTURE_WITH_CAPTION = 8
PICTURE_PLACEHOLDER_IDX = 1


@pytest.fixture(scope="module")
def pdf_deck(decks_dir: Path) -> Deck:
    return ingest_slides(decks_dir / "lecture01.pdf")


@pytest.fixture(scope="module")
def pptx_deck(decks_dir: Path) -> Deck:
    return ingest_slides(decks_dir / "lecture01.pptx")


# --- the fixture -------------------------------------------------------------------


def test_pdf_slide_3_yields_one_240_by_150_png_and_slides_1_and_2_none(pdf_deck: Deck) -> None:
    (asset,) = pdf_deck.assets
    assert asset.media_type == "image/png"
    assert (asset.width, asset.height) == FIGURE_SIZE
    assert [slide.image_ids for slide in pdf_deck.slides] == [(), (), (asset.id,)]
    assert pdf_deck.recurring_image_ids == ()


def test_pdf_figure_is_reencoded_so_its_id_differs_from_the_pptx_one(pdf_deck: Deck) -> None:
    """pypdf hands back the Flate stream re-encoded as PNG; ids hash bytes as extracted
    (ticket decision: nothing needs cross-format identity - a lecture has one deck)."""
    (asset,) = pdf_deck.assets
    assert asset.id != PPTX_IMAGE_ID
    assert asset.id == image_id(asset.data)


def test_pptx_slide_3_asset_id_is_the_committed_png_hash(pptx_deck: Deck, decks_dir: Path) -> None:
    (asset,) = pptx_deck.assets
    assert asset.id == PPTX_IMAGE_ID == image_id((decks_dir / "value_iteration.png").read_bytes())
    assert pptx_deck.slides[2].image_ids == (PPTX_IMAGE_ID,)


# --- ad-hoc PPTX decks -------------------------------------------------------------


def _png(tmp_path: Path, name: str, size: tuple[int, int], colour: str) -> Path:
    path = tmp_path / f"{name}.png"
    PILImage.new("RGB", size, colour).save(path)
    return path


def _save(prs: PresentationType, tmp_path: Path, name: str = "adhoc.pptx") -> Path:
    path = tmp_path / name
    prs.save(str(path))
    return path


def _blank_slide(prs: PresentationType) -> SlideShapes:
    return prs.slides.add_slide(prs.slide_layouts[BLANK]).shapes


def _deck_with_pictures(tmp_path: Path, per_slide: list[list[Path]]) -> Path:
    """One blank slide per entry; its pictures side by side, left to right."""
    prs = Presentation()
    for pictures in per_slide:
        shapes = _blank_slide(prs)
        for i, picture in enumerate(pictures):
            shapes.add_picture(str(picture), Inches(1 + 3 * i), Inches(1))
    return _save(prs, tmp_path)


@pytest.fixture
def figure(tmp_path: Path) -> Path:
    return _png(tmp_path, "figure", (64, 48), "navy")


@pytest.fixture
def other_figure(tmp_path: Path) -> Path:
    return _png(tmp_path, "other", (64, 48), "olive")


@pytest.fixture
def tiny(tmp_path: Path) -> Path:
    return _png(tmp_path, "tiny", (8, 8), "red")


def test_tiny_picture_is_dropped_from_image_ids_and_assets_by_default(
    tmp_path: Path, tiny: Path, figure: Path
) -> None:
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[tiny, figure]]))
    (only,) = deck.slides
    assert only.image_ids == (image_id(figure.read_bytes()),)
    assert [asset.id for asset in deck.assets] == [image_id(figure.read_bytes())]


def test_min_px_4_keeps_the_tiny_picture(tmp_path: Path, tiny: Path) -> None:
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[tiny]]), min_px=4)
    (only,) = deck.slides
    assert only.image_ids == (image_id(tiny.read_bytes()),)
    assert (deck.assets[0].width, deck.assets[0].height) == (8, 8)


def test_min_px_applies_to_either_dimension(tmp_path: Path) -> None:
    """A 200 x 8 rule is a decoration even though it is wide."""
    rule = _png(tmp_path, "rule", (200, 8), "gray")
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[rule]]))
    assert deck.slides[0].image_ids == () and deck.assets == ()


def test_picture_on_3_of_4_slides_is_recurring_and_leaves_every_image_ids(
    tmp_path: Path, figure: Path, other_figure: Path
) -> None:
    logo, other = image_id(figure.read_bytes()), image_id(other_figure.read_bytes())
    deck = ingest_slides(
        _deck_with_pictures(tmp_path, [[figure], [figure, other_figure], [figure], []])
    )
    assert deck.recurring_image_ids == (logo,)
    assert [asset.id for asset in deck.assets].count(logo) == 1
    assert [slide.image_ids for slide in deck.slides] == [(), (other,), (), ()]


def test_picture_on_2_of_4_slides_is_not_recurring(tmp_path: Path, figure: Path) -> None:
    """The rule is *more than half*: two of four is exactly half."""
    logo = image_id(figure.read_bytes())
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[figure], [figure], [], []]))
    assert deck.recurring_image_ids == ()
    assert [slide.image_ids for slide in deck.slides] == [(logo,), (logo,), (), ()]


def test_two_slide_deck_with_the_same_picture_on_both_is_not_recurring(
    tmp_path: Path, figure: Path
) -> None:
    """Two slides cannot tell a logo from a figure shown twice; the rule needs >= 3."""
    logo = image_id(figure.read_bytes())
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[figure], [figure]]))
    assert deck.recurring_image_ids == ()
    assert [slide.image_ids for slide in deck.slides] == [(logo,), (logo,)]


def test_picture_and_text_box_inside_a_group_are_found_in_the_groups_slot(
    tmp_path: Path, figure: Path
) -> None:
    prs = Presentation()
    shapes = _blank_slide(prs)
    below = shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
    below.text_frame.text = "below the group"
    group = shapes.add_group_shape()
    group.shapes.add_picture(str(figure), Inches(1), Inches(1))
    inside = group.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    inside.text_frame.text = "inside the group"
    (only,) = ingest_slides(_save(prs, tmp_path)).slides
    assert only.image_ids == (image_id(figure.read_bytes()),)
    assert only.blocks == (
        TextBlock(lines=("inside the group",)),
        TextBlock(lines=("below the group",)),
    )


def test_two_different_pictures_on_one_slide_are_listed_left_to_right(
    tmp_path: Path, figure: Path, other_figure: Path
) -> None:
    prs = Presentation()
    shapes = _blank_slide(prs)
    shapes.add_picture(str(other_figure), Inches(5), Inches(1))  # added first, placed right
    shapes.add_picture(str(figure), Inches(1), Inches(1))
    (only,) = ingest_slides(_save(prs, tmp_path)).slides
    assert only.image_ids == (image_id(figure.read_bytes()), image_id(other_figure.read_bytes()))


def test_same_picture_twice_on_one_slide_is_one_asset_listed_once(
    tmp_path: Path, figure: Path
) -> None:
    deck = ingest_slides(_deck_with_pictures(tmp_path, [[figure, figure]]))
    assert len(deck.assets) == 1
    assert deck.slides[0].image_ids == (image_id(figure.read_bytes()),)


def test_picture_dropped_into_a_picture_placeholder_counts(tmp_path: Path, figure: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[PICTURE_WITH_CAPTION])
    slide.shapes.title.text = "Captioned"
    slide.placeholders[PICTURE_PLACEHOLDER_IDX].insert_picture(str(figure))
    (only,) = ingest_slides(_save(prs, tmp_path)).slides
    assert only.title == "Captioned"
    assert only.image_ids == (image_id(figure.read_bytes()),)


def test_parsers_stay_faithful_and_the_rules_live_in_ingest_slides(
    tmp_path: Path, tiny: Path, figure: Path
) -> None:
    """``parse_pptx`` reports the tiny picture and the logo; ``ingest_slides`` applies
    the rules on top - the same structure as ``merge_sentences`` behind
    ``ingest_captions``."""
    path = _deck_with_pictures(tmp_path, [[tiny, figure], [figure], [figure]])
    raw = parse_pptx(path)
    assert len(raw.assets) == 2 and raw.recurring_image_ids == ()
    assert raw.slides[0].image_ids == (image_id(tiny.read_bytes()), image_id(figure.read_bytes()))
    cooked = ingest_slides(path)
    assert cooked == set_aside_recurring_images(drop_small_images(raw, min_px=32))
    assert cooked.recurring_image_ids == (image_id(figure.read_bytes()),)
    assert all(slide.image_ids == () for slide in cooked.slides)


# --- the rules on hand-built decks -------------------------------------------------


def _asset(seed: bytes, width: int = 100, height: int = 100) -> SlideImage:
    return SlideImage(
        id=image_id(seed), media_type="image/png", width=width, height=height, data=seed
    )


def _slide(number: int, *image_ids: str) -> Slide:
    return Slide(number=number, title=None, blocks=(), notes=None, image_ids=image_ids)


def test_drop_small_images_removes_the_asset_and_every_reference() -> None:
    small, big = _asset(b"small", 8, 8), _asset(b"big")
    deck = Deck(
        source="x.pptx",
        slides=(_slide(1, small.id, big.id), _slide(2, small.id)),
        assets=(small, big),
    )
    kept = drop_small_images(deck, min_px=32)
    assert kept.assets == (big,)
    assert [slide.image_ids for slide in kept.slides] == [(big.id,), ()]
    assert drop_small_images(deck, min_px=8) == deck


def test_set_aside_recurring_images_lists_ids_in_first_seen_order() -> None:
    a, b, c = _asset(b"a"), _asset(b"b"), _asset(b"c")
    deck = Deck(
        source="x.pptx",
        slides=(_slide(1, b.id, a.id), _slide(2, a.id, b.id, c.id), _slide(3, a.id, b.id)),
        assets=(a, b, c),
    )
    result = set_aside_recurring_images(deck)
    assert result.recurring_image_ids == (b.id, a.id)
    assert result.assets == (a, b, c)
    assert [slide.image_ids for slide in result.slides] == [(), (c.id,), ()]
    assert set_aside_recurring_images(result) == result


# --- media types -------------------------------------------------------------------


@pytest.mark.parametrize(
    ("media_type", "kept"),
    [
        ("image/png", True),
        ("image/jpeg", True),
        ("image/gif", True),
        ("image/bmp", True),
        ("image/tiff", True),
        ("image/webp", True),
        ("image/x-emf", False),
        ("image/x-wmf", False),
        ("image/svg+xml", False),
    ],
)
def test_media_type_table(media_type: str, kept: bool) -> None:
    assert _media_type_ok(media_type) is kept


# --- broken PDF images -------------------------------------------------------------


def _page_3_with_broken_image(
    decks_dir: Path, tmp_path: Path, corrupt: Callable[[DictionaryObject], None]
) -> Path:
    writer = PdfWriter()
    writer.add_page(PdfReader(str(decks_dir / "lecture01.pdf")).pages[2])
    xobjects = writer.pages[0]["/Resources"]["/XObject"]
    for name in xobjects:
        corrupt(xobjects[name].get_object())
    path = tmp_path / "broken.pdf"
    writer.write(str(path))
    return path


def _truncate_stream(image: DictionaryObject) -> None:
    image._data = b"\x00" * 7  # type: ignore[attr-defined]


def _claim_jpeg(image: DictionaryObject) -> None:
    image[NameObject("/Filter")] = NameObject("/DCTDecode")
    image._data = b"\xff\xd8not a jpeg"  # type: ignore[attr-defined]


@pytest.mark.parametrize("corrupt", [_truncate_stream, _claim_jpeg], ids=["pypdf", "pillow"])
def test_broken_pdf_image_is_skipped_and_the_slide_text_kept(
    decks_dir: Path, tmp_path: Path, corrupt: Callable[[DictionaryObject], None]
) -> None:
    """Losing one picture is recoverable; raising would lose the whole deck."""
    (only,) = parse_pdf(_page_3_with_broken_image(decks_dir, tmp_path, corrupt)).slides
    assert only.title == "Value Iteration"
    assert len(only.blocks) == 1
    assert only.image_ids == ()
