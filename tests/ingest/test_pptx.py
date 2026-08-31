"""P2-01: ``parse_pptx`` on the fixture deck, ad-hoc decks, and ``clean_line``.

Each row of the decks table in ``tests/fixtures/README.md`` is a test name here. The
ad-hoc cases are built in-memory with python-pptx and saved under ``tmp_path`` - the
slide-side analogue of Phase 1's inline VTT strings - because they exercise structural
edge cases of the file format, not the lecture.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches

from lecturenotes.ingest.slides import (
    Deck,
    DeckParseError,
    TextBlock,
    clean_line,
    image_id,
    parse_pptx,
)

# README decks table, transcribed - never imported from make_deck.py (which is
# fixture-group-only code) nor read from the parser output.
TITLES = ("Markov Decision Processes", "The Bellman Equation", "Value Iteration")
IMAGE_ID = "img-a63ae9b7dc5e9397"
EXAM_PHRASE = "this will be on the exam"

TITLE_AND_CONTENT = 1
BLANK = 6


@pytest.fixture(scope="module")
def deck(decks_dir: Path) -> Deck:
    return parse_pptx(decks_dir / "lecture01.pptx")


@pytest.fixture(scope="module")
def png(decks_dir: Path) -> bytes:
    return (decks_dir / "value_iteration.png").read_bytes()


# --- the decks table ---------------------------------------------------------------


def test_slide_1_single_column_gives_title_then_five_bullets(deck: Deck) -> None:
    slide = deck.slides[0]
    assert slide.title == TITLES[0]
    assert len(slide.blocks) == 1
    assert len(slide.blocks[0].lines) == 5
    assert slide.blocks[0].lines[0].startswith("States s in S")


def test_slide_2_two_columns_read_left_block_then_right_block(deck: Deck) -> None:
    slide = deck.slides[1]
    assert len(slide.blocks) == 2
    left, right = slide.blocks
    assert left.lines[0] == "Equation" and len(left.lines) == 6
    assert right.lines[0] == "Intuition" and len(right.lines) == 6


def test_slide_3_steps_in_order_and_one_png_figure(deck: Deck, png: bytes) -> None:
    slide = deck.slides[2]
    assert len(slide.blocks) == 1
    assert [line[:2] for line in slide.blocks[0].lines] == ["1.", "2.", "ga", "3.", "4."]
    assert slide.image_ids == (IMAGE_ID,)
    (asset,) = deck.assets
    assert asset.id == IMAGE_ID == image_id(png)
    assert asset.media_type == "image/png"
    assert (asset.width, asset.height) == (240, 150)
    assert asset.data == png


def test_every_slide_has_speaker_notes_and_slide_2_flags_the_exam(deck: Deck) -> None:
    assert all(slide.notes is not None for slide in deck.slides)
    assert deck.slides[1].notes is not None and EXAM_PHRASE in deck.slides[1].notes


def test_titles_come_from_the_title_placeholder(deck: Deck) -> None:
    assert tuple(slide.title for slide in deck.slides) == TITLES


def test_slide_numbers_are_file_positions_and_nothing_is_hidden(deck: Deck) -> None:
    assert [slide.number for slide in deck.slides] == [1, 2, 3]
    assert not any(slide.hidden for slide in deck.slides)
    assert deck.recurring_image_ids == ()


def test_source_is_the_path_as_given(decks_dir: Path) -> None:
    path = decks_dir / "lecture01.pptx"
    assert parse_pptx(path).source == str(path)


# --- ad-hoc decks ------------------------------------------------------------------


def _save(prs: PresentationType, tmp_path: Path, name: str = "adhoc.pptx") -> Path:
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_empty_body_placeholder_yields_no_block(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "Only a title"
    (only,) = parse_pptx(_save(prs, tmp_path)).slides
    assert only.title == "Only a title"
    assert only.blocks == ()


def test_blank_layout_text_box_has_no_title_and_one_block(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "free text"
    (only,) = parse_pptx(_save(prs, tmp_path)).slides
    assert only.title is None
    assert only.blocks == (TextBlock(lines=("free text",)),)


def test_soft_line_break_stays_within_one_line(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.add_run().text = "a wrapped"
    paragraph.add_line_break()
    paragraph.add_run().text = "bullet"
    box.text_frame.add_paragraph().text = "next bullet"
    assert "\v" in box.text_frame.text  # python-pptx reads a soft break back as \v
    (only,) = parse_pptx(_save(prs, tmp_path)).slides
    assert only.blocks == (TextBlock(lines=("a wrapped bullet", "next bullet")),)


def test_table_becomes_one_line_per_row_with_cells_joined(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
    table.cell(0, 0).text, table.cell(0, 1).text = "State", "Value"
    table.cell(1, 0).text, table.cell(1, 1).text = "s1", "0.5"
    (only,) = parse_pptx(_save(prs, tmp_path)).slides
    assert only.blocks == (TextBlock(lines=("State | Value", "s1 | 0.5")),)


def test_hidden_slide_keeps_its_number_and_is_flagged(tmp_path: Path) -> None:
    prs = Presentation()
    for i in (1, 2, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
        slide.shapes.title.text = f"Slide {i}"
        if i == 2:
            slide._element.set("show", "0")
    slides = parse_pptx(_save(prs, tmp_path)).slides
    assert [(s.number, s.hidden) for s in slides] == [(1, False), (2, True), (3, False)]


def test_blank_or_absent_notes_give_none(tmp_path: Path) -> None:
    prs = Presentation()
    blank_notes = prs.slides.add_slide(prs.slide_layouts[BLANK])
    blank_notes.notes_slide.notes_text_frame.text = " \t "
    no_notes = prs.slides.add_slide(prs.slide_layouts[BLANK])
    assert not no_notes.has_notes_slide
    slides = parse_pptx(_save(prs, tmp_path)).slides
    assert [s.notes for s in slides] == [None, None]


def test_text_box_above_the_body_is_read_first_even_when_added_last(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = "Title"
    slide.placeholders[1].text_frame.text = "body"
    above = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
    above.text_frame.text = "above the body"
    (only,) = parse_pptx(_save(prs, tmp_path)).slides
    assert only.blocks == (TextBlock(lines=("above the body",)), TextBlock(lines=("body",)))


def test_garbage_bytes_raise_deck_parse_error_naming_the_file(tmp_path: Path) -> None:
    junk = tmp_path / "junk.pptx"
    junk.write_bytes(b"\x00not a zip\xff" * 64)
    with pytest.raises(DeckParseError, match="junk.pptx"):
        parse_pptx(junk)


# --- clean_line --------------------------------------------------------------------


@pytest.mark.parametrize(
    ("raw", "expected"),
    [
        ("- States s in S", "States s in S"),
        ("• x", "x"),  # bullet
        ("– x", "x"),  # en dash
        ("— x", "x"),  # em dash
        ("* x", "x"),
        ("▪ x", "x"),  # small black square
        ("1. Initialise", "1. Initialise"),  # numbered markers are content
        ("-x", "-x"),  # the glyph must be followed by whitespace
        ("a\vb", "a b"),
        ("a\u00a0b\tc", "a b c"),  # NBSP and tab collapse to one space
        ("  padded  ", "padded"),
        ("-  -  twice", "- twice"),  # exactly one glyph is removed
        ("", ""),
        ("- ", ""),
    ],
)
def test_clean_line_table(raw: str, expected: str) -> None:
    assert clean_line(raw) == expected


def test_clean_line_is_idempotent_and_identity_on_every_fixture_line(deck: Deck) -> None:
    texts: list[str] = []
    for slide in deck.slides:
        texts.extend(t for t in (slide.title, slide.notes) if t is not None)
        texts.extend(line for block in slide.blocks for line in block.lines)
    assert texts
    for text in texts:
        assert clean_line(text) == text
