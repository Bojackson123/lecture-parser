"""Slide ingest (plan §3 stage 2): ``.pptx`` / ``.pdf`` → ``Deck``.

Two parsers, one pure layout function and a composing entrypoint, added one ticket at
a time:

    Deck / Slide / TextBlock / SlideImage, clean_line   P2-01
    parse_pptx(path)       → Deck                       P2-01  (text + notes + pictures)
    ingest_slides(path)    → Deck                       P2-01  (dispatch by suffix)
    layout_page(spans)     → PageLayout                 P2-02  (pure; columns, title)
    parse_pdf(path)        → Deck                       P2-02  (text; boilerplate dropped)
    image rules            in ingest_slides             P2-03  (size, recurring; groups)

The stage is pure: image bytes stay in ``SlideImage.data`` and nothing is written to
disk. Phase 5 mints ``MediaAsset``s from ``SlideImage``s and owns where the bytes go.
Both parsers push every title, body line and note through ``clean_line`` so the PDF's
``- States…`` equals the PPTX's ``States…`` (P2-02's cross-format invariant).
"""

from __future__ import annotations

import hashlib
import math
import re
import zipfile
from collections import Counter
from collections.abc import Callable, Iterable, Iterator, Sequence
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu, Length
from pydantic import BaseModel, ConfigDict, model_validator
from pypdf import PageObject, PasswordType, PdfReader
from pypdf.errors import DependencyError, PyPdfError

__all__ = [
    "Deck",
    "DeckParseError",
    "PageLayout",
    "Slide",
    "SlideImage",
    "Span",
    "TextBlock",
    "clean_line",
    "drop_small_images",
    "image_id",
    "ingest_slides",
    "layout_page",
    "parse_pdf",
    "parse_pptx",
    "set_aside_recurring_images",
]


class _DeckModel(BaseModel):
    # Bytes travel as base64 in JSON, so a whole ``Deck`` is one plain JSON document:
    # the expected fixture is a file, and P2-04's ``--json`` output can be re-read.
    model_config = ConfigDict(
        frozen=True, extra="forbid", ser_json_bytes="base64", val_json_bytes="base64"
    )


def _duplicates(ids: Iterable[str]) -> list[str]:
    return sorted(value for value, n in Counter(ids).items() if n > 1)


def image_id(data: bytes) -> str:
    """``img-`` + the first 16 hex digits of ``sha256(data)``.

    Content-addressed, not ``slideN-imgM``: a figure reused on two slides is one asset
    with two references, and inserting a slide moves nothing (cf. plan §7.2).
    """
    return "img-" + hashlib.sha256(data).hexdigest()[:16]


class TextBlock(_DeckModel):
    """One shape's worth of text: one line per paragraph (or per table row)."""

    lines: tuple[str, ...]

    @model_validator(mode="after")
    def _valid(self) -> TextBlock:
        if not self.lines:
            raise ValueError("a text block needs at least one line")
        if any(not line for line in self.lines):
            raise ValueError("text block lines must not be empty")
        return self


class SlideImage(_DeckModel):
    """A picture's bytes, identified by their hash (``image_id``)."""

    id: str
    media_type: str
    width: int
    height: int
    data: bytes

    @model_validator(mode="after")
    def _valid(self) -> SlideImage:
        expected = image_id(self.data)
        if self.id != expected:
            raise ValueError(f"image id {self.id!r} does not match its content hash {expected!r}")
        if self.width < 1 or self.height < 1:
            raise ValueError(f"width and height must be >= 1, got {self.width}x{self.height}")
        return self


class Slide(_DeckModel):
    """One slide. ``number`` is the 1-based position in the file, hidden slides included.

    That is what a reader counts when they open the deck, so a ``SlideRange`` in an
    anchor stays honest; Phase 4 may skip ``hidden`` slides but must not renumber.
    """

    number: int
    title: str | None
    blocks: tuple[TextBlock, ...]
    notes: str | None
    image_ids: tuple[str, ...]
    hidden: bool = False

    @model_validator(mode="after")
    def _valid(self) -> Slide:
        if self.number < 1:
            raise ValueError(f"slide number must be >= 1, got {self.number}")
        dupes = _duplicates(self.image_ids)
        if dupes:
            raise ValueError(f"duplicate image id(s) on slide {self.number}: {', '.join(dupes)}")
        return self


class Deck(_DeckModel):
    """The output of stage 2.

    ``source`` is the path as given, in POSIX form (for ``SourceRef.deck_path``).
    """

    source: str
    slides: tuple[Slide, ...]
    assets: tuple[SlideImage, ...]
    recurring_image_ids: tuple[str, ...] = ()

    @model_validator(mode="after")
    def _valid(self) -> Deck:
        numbers = [slide.number for slide in self.slides]
        if numbers != list(range(1, len(numbers) + 1)):
            raise ValueError(
                f"slide numbers must be exactly 1..{len(numbers)} in order, got {numbers}"
            )
        dupes = _duplicates(asset.id for asset in self.assets)
        if dupes:
            raise ValueError(f"duplicate asset id(s): {', '.join(dupes)}")
        known = {asset.id for asset in self.assets}
        missing = sorted({i for slide in self.slides for i in slide.image_ids if i not in known})
        if missing:
            raise ValueError(f"slide(s) reference unknown image id(s): {', '.join(missing)}")
        missing = sorted({i for i in self.recurring_image_ids if i not in known})
        if missing:
            raise ValueError(
                f"recurring_image_ids reference unknown image id(s): {', '.join(missing)}"
            )
        return self


class DeckParseError(ValueError):
    """A deck file that could not be opened or read; names the file and the cause."""

    def __init__(self, path: Path, cause: BaseException) -> None:
        super().__init__(f"{path}: {cause}")
        self.path = path
        self.__cause__ = cause


# --- line cleaning -----------------------------------------------------------------

_WHITESPACE = re.compile(r"\s+")  # \s covers \v, \t and NBSP
_LEADING_BULLET = re.compile(r"^[-–—•·▪●○■‣*]\s")


def clean_line(text: str) -> str:
    """Collapse whitespace, then drop exactly one leading bullet glyph.

    A soft line break (``\\v``) and NBSP become a space; a bullet glyph is removed only
    when followed by whitespace, so ``-x`` and ``1. Initialise`` are content. One glyph,
    not many: ``-  -  twice`` is a line that starts with a dash. Idempotent on cleaned
    text and the identity on clean text.
    """
    text = _WHITESPACE.sub(" ", text).strip()
    return _LEADING_BULLET.sub("", text, count=1).strip()


def _clean_or_none(text: str | None) -> str | None:
    if text is None:
        return None
    cleaned = clean_line(text)
    return cleaned or None


# --- images: the shared allow-list ---------------------------------------------------
#
# Raster formats every downstream target can inline. Vector images (EMF/WMF/SVG) are
# skipped, not converted: rasterising needs a system renderer, nothing downstream can
# inline them, and they are almost always decorations (P2-03 decision; a deck that keeps
# its figures as EMF is a Phase 9 concern - render the slide).

_MEDIA_TYPES = frozenset(
    {"image/png", "image/jpeg", "image/gif", "image/bmp", "image/tiff", "image/webp"}
)
_MIME_BY_EXT = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".tif": "image/tiff",
    ".tiff": "image/tiff",
    ".webp": "image/webp",
}


def _media_type_ok(media_type: str) -> bool:
    return media_type in _MEDIA_TYPES


def _slide_image(data: bytes, media_type: str, size: tuple[int, int]) -> SlideImage:
    width, height = size
    return SlideImage(
        id=image_id(data), media_type=media_type, width=width, height=height, data=data
    )


# --- PPTX -----------------------------------------------------------------------------

_TITLE_TYPES = frozenset({PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE})
_DEFAULT_SLIDE_HEIGHT = Emu(6858000)  # 7.5 in, python-pptx's default template
_ROW_BANDS = 20


def _is_title(shape: BaseShape) -> bool:
    return shape.is_placeholder and shape.placeholder_format.type in _TITLE_TYPES


def _reading_order(shapes: Iterable[BaseShape], slide_height: Length) -> list[BaseShape]:
    """Sibling shapes sorted by (row band, left) — position, not z-order.

    spTree order is authoring order, so a text box added last but placed at the top
    must still come first. Bands of ¹⁄₂₀ slide height keep two side-by-side
    placeholders (same top) in left-to-right order.
    """
    band = max(int(slide_height) // _ROW_BANDS, 1)

    def key(shape: BaseShape) -> tuple[int, int]:
        return (int(shape.top or 0) // band, int(shape.left or 0))

    return sorted(shapes, key=key)


def _walk(shapes: Iterable[BaseShape], slide_height: Length) -> Iterator[BaseShape]:
    """Leaf shapes in reading order, groups flattened in place.

    A ``GroupShape`` takes the positional slot of its own ``top``/``left``; its children
    follow there, in their own reading order (child coordinates share one space, so
    the same key sorts them). Groups nest, so this recurses.
    """
    for shape in _reading_order(shapes, slide_height):
        if isinstance(shape, GroupShape):
            yield from _walk(shape.shapes, slide_height)
        else:
            yield shape


def _picture(shape: Picture) -> SlideImage | None:
    """The picture's ``SlideImage``; ``None`` when it is skipped.

    Skipped: a media type outside the allow-list (python-pptx reports EMF/WMF as
    ``image/x-wmf``), and any picture whose bytes Pillow cannot identify or whose part
    is missing — python-pptx raises from ``.image`` in those cases, and a broken picture
    must never lose a slide's text.
    """
    try:
        image = shape.image
        media_type = image.content_type
        if not _media_type_ok(media_type):
            return None
        return _slide_image(image.blob, media_type, image.size)
    except (OSError, ValueError, KeyError):  # UnidentifiedImageError is an OSError
        return None


def _shape_lines(shape: BaseShape) -> list[str]:
    """Cleaned, non-empty lines of a text-bearing or table shape; ``[]`` otherwise."""
    if isinstance(shape, GraphicFrame) and shape.has_table:
        raw = [" | ".join(cell.text for cell in row.cells) for row in shape.table.rows]
    elif shape.has_text_frame:
        # ``has_text_frame`` is only true on shapes that carry a ``text_frame``.
        raw = [paragraph.text for paragraph in shape.text_frame.paragraphs]  # type: ignore[attr-defined]
    else:
        return []
    return [line for line in (clean_line(r) for r in raw) if line]


def _notes(slide: PptxSlide) -> str | None:
    if not slide.has_notes_slide:
        return None
    frame = slide.notes_slide.notes_text_frame
    return _clean_or_none(frame.text if frame is not None else None)


def parse_pptx(path: Path) -> Deck:
    """Read a ``.pptx`` into a ``Deck``: titles, body blocks, notes, pictures.

    Per slide, ``title`` is the ``TITLE``/``CENTER_TITLE`` placeholder; every other
    shape with a text frame or a table becomes one ``TextBlock`` (a line per paragraph,
    or per table row with cells joined by ``" | "``), in reading order; each picture
    (a picture dropped into a placeholder included) becomes a ``SlideImage`` listed once
    per slide and once in ``assets``. Group shapes are flattened in place (``_walk``).
    Every raster picture is reported: the size and recurring rules are
    ``ingest_slides``'s, so the parser stays a faithful extractor. Raises
    ``DeckParseError`` for a file that is not a presentation; lets
    ``FileNotFoundError`` through.
    """
    if not path.is_file():
        raise FileNotFoundError(path)
    try:
        prs = Presentation(str(path))
    except (PackageNotFoundError, zipfile.BadZipFile, KeyError) as exc:
        raise DeckParseError(path, exc) from exc
    slide_height = prs.slide_height if prs.slide_height is not None else _DEFAULT_SLIDE_HEIGHT

    slides: list[Slide] = []
    assets: dict[str, SlideImage] = {}
    for number, pptx_slide in enumerate(prs.slides, start=1):
        title: str | None = None
        blocks: list[TextBlock] = []
        image_ids: list[str] = []
        for shape in _walk(pptx_slide.shapes, slide_height):
            if _is_title(shape):
                if title is None:
                    title = _clean_or_none(shape.text_frame.text)  # type: ignore[attr-defined]
                continue
            if isinstance(shape, Picture):
                asset = _picture(shape)
                if asset is not None:
                    assets.setdefault(asset.id, asset)
                    if asset.id not in image_ids:
                        image_ids.append(asset.id)
                continue
            lines = _shape_lines(shape)
            if lines:
                blocks.append(TextBlock(lines=tuple(lines)))
        slides.append(
            Slide(
                number=number,
                title=title,
                blocks=tuple(blocks),
                notes=_notes(pptx_slide),
                image_ids=tuple(image_ids),
                hidden=pptx_slide._element.get("show") == "0",
            )
        )
    # ``as_posix``: the path as given, with forward slashes on every platform, so the
    # expected-deck fixture (and any ``--json`` output) is byte-identical across OSes.
    return Deck(source=path.as_posix(), slides=tuple(slides), assets=tuple(assets.values()))


# --- PDF: page layout (pure) -------------------------------------------------------------
#
# A PDF has no placeholders, only positioned strings. ``layout_page`` turns one page's
# spans into a title plus one ``TextBlock`` per column, by rules that never look at the
# order the spans arrived in (the content-stream order is what must not leak out):
#
#   rows     spans whose baselines differ by <= 0.5 x min(size) share a row
#   title    the topmost row set entirely in the largest face, if that face beats every
#            other size by 1.15x (or the page is that one row); a same-size row directly
#            beneath (gap <= 1.5 x size) joins it
#   columns  single linkage over the distinct x-starts: a gap wider than 0.15 x page
#            width starts a new column, so a sub-bullet indent never does
#   order    columns left to right; rows top to bottom; spans in a row joined by a space

_ROW_TOLERANCE = 0.5
_TITLE_RATIO = 1.15
_TITLE_JOIN = 1.5
_COLUMN_GAP = 0.15


class Span(_DeckModel):
    """One drawn string: start point in PDF user space (``y`` grows upward) and size."""

    x: float
    y: float
    size: float
    text: str

    @model_validator(mode="after")
    def _valid(self) -> Span:
        if not self.size > 0:
            raise ValueError(f"span size must be > 0, got {self.size}")
        return self


class PageLayout(_DeckModel):
    """``layout_page``'s result: the page title, if any, and one ``TextBlock`` per column."""

    title: str | None
    blocks: tuple[TextBlock, ...]


_Row = list[Span]  # spans sharing a baseline, left to right


def _top_to_bottom(spans: Iterable[Span]) -> list[Span]:
    # A total order (every field takes part) so equal multisets of spans sort equally.
    return sorted(spans, key=lambda s: (-s.y, s.x, s.size, s.text))


def _rows(spans: Iterable[Span]) -> list[_Row]:
    """Group spans into rows top to bottom, each row left to right."""
    rows: list[_Row] = []
    for span in _top_to_bottom(spans):
        if rows:
            above = rows[-1][-1]
            if above.y - span.y <= _ROW_TOLERANCE * min(above.size, span.size):
                rows[-1].append(span)
                continue
        rows.append([span])
    return [sorted(row, key=lambda s: (s.x, -s.y, s.size, s.text)) for row in rows]


def _row_y(row: _Row) -> float:
    return max(span.y for span in row)


def _row_text(row: _Row) -> str:
    return clean_line(" ".join(span.text for span in row))


def _split_title(rows: list[_Row]) -> tuple[str | None, list[_Row]]:
    """The title (one or two rows in the largest face) and the rows that remain."""
    top = max(span.size for row in rows for span in row)
    starts = [i for i, row in enumerate(rows) if all(span.size == top for span in row)]
    if not starts:
        return None, rows
    start = starts[0]
    end = start + 1
    if (
        end in starts
        and _row_y(rows[start]) - _row_y(rows[end]) <= _TITLE_JOIN * top
    ):
        end += 1
    rest = rows[:start] + rows[end:]
    others = [span.size for row in rest for span in row]
    if others:
        if top < _TITLE_RATIO * max(others):
            return None, rows
    elif len(rows) != 1:
        return None, rows
    return clean_line(" ".join(_row_text(row) for row in rows[start:end])), rest


def _columns(spans: Iterable[Span], page_width: float) -> list[list[Span]]:
    """Single-linkage clusters of x-starts, left to right."""
    gap = _COLUMN_GAP * page_width
    columns: list[list[Span]] = []
    previous_x: float | None = None
    for span in sorted(spans, key=lambda s: (s.x, -s.y, s.size, s.text)):
        if previous_x is None or span.x - previous_x > gap:
            columns.append([])
        columns[-1].append(span)
        previous_x = span.x
    return columns


def layout_page(spans: Iterable[Span], *, page_width: float, page_height: float) -> PageLayout:
    """Title and column blocks of one page, independent of the order of ``spans``.

    Spans whose ``clean_line`` text is empty are ignored. Every emitted line goes
    through ``clean_line``. ``page_height`` is accepted with the page box for
    symmetry; no rule depends on it (boilerplate is found by cross-page repetition in
    ``parse_pdf``, not by margin position).
    """
    kept = [span for span in spans if clean_line(span.text)]
    if not kept:
        return PageLayout(title=None, blocks=())
    title, body_rows = _split_title(_rows(kept))
    blocks: list[TextBlock] = []
    for column in _columns((span for row in body_rows for span in row), page_width):
        lines = [text for text in (_row_text(row) for row in _rows(column)) if text]
        if lines:
            blocks.append(TextBlock(lines=tuple(lines)))
    return PageLayout(title=title, blocks=tuple(blocks))


# --- PDF: the file -------------------------------------------------------------------------

_DIGITS = re.compile(r"\d+")
_Matrix = Sequence[float]  # a PDF affine matrix [a b c d e f]


def _compose(tm: _Matrix, cm: _Matrix) -> tuple[float, float, float]:
    """Text-space origin and unit height of ``tm`` x ``cm``: ``(x, y, scale)``.

    pypdf hands the visitor the raw text matrix and the current transformation
    separately; exporters that draw in a scaled ``cm`` (PowerPoint, Cairo) put the
    real size in the matrices rather than in ``Tf``.
    """
    x = tm[4] * cm[0] + tm[5] * cm[2] + cm[4]
    y = tm[4] * cm[1] + tm[5] * cm[3] + cm[5]
    c = tm[2] * cm[0] + tm[3] * cm[2]
    d = tm[2] * cm[1] + tm[3] * cm[3]
    return x, y, math.hypot(c, d)


def _page_spans(page: PageObject) -> list[Span]:
    spans: list[Span] = []

    def visit(text: str, cm: _Matrix, tm: _Matrix, font: object, font_size: float) -> None:
        text = text.rstrip("\n")  # pypdf appends one at every line break it detects
        if not text.strip():
            return
        x, y, scale = _compose(tm, cm)
        size = float(font_size) * scale
        spans.append(Span(x=x, y=y, size=size if size > 0 else float(font_size), text=text))

    page.extract_text(visitor_text=visit)
    return spans


def _page_images(page: PageObject) -> list[SlideImage]:
    """The page's raster images as pypdf yields them, unknown and broken ones skipped.

    ``page.images[i]`` decodes lazily and raises on some filters (pypdf) and some
    streams (Pillow), so images are fetched one at a time and a failure skips **that
    image** with no error: losing one picture is recoverable, raising would lose the
    slide's text and, with ``ingest_slides`` as the only entrypoint, the whole deck.
    Nothing is logged - the package has no logging yet; P2-04's command shows what was
    found.
    """
    images = page.images
    found: list[SlideImage] = []
    for i in range(len(images)):
        try:
            img = images[i]
            media_type = _MIME_BY_EXT.get(Path(img.name).suffix.lower())
            if media_type is None or not _media_type_ok(media_type) or img.image is None:
                continue
            found.append(_slide_image(img.data, media_type, img.image.size))
        except Exception:  # any decode failure, deliberately broad; see the docstring
            continue
    return found


def _boilerplate(pages: Sequence[Sequence[Span]]) -> set[str]:
    """Digit-normalised span texts on more than half the pages of a deck of >= 2."""
    if len(pages) < 2:
        return set()
    counts = Counter(key for page in pages for key in {_normalised(span) for span in page})
    return {key for key, n in counts.items() if n > len(pages) / 2}


def _normalised(span: Span) -> str:
    return _DIGITS.sub("#", clean_line(span.text))


def parse_pdf(path: Path) -> Deck:
    """Read a ``.pdf`` into a ``Deck``: one slide per page, title and column blocks.

    Running headers and footers - any span whose digit-normalised text recurs on more
    than half the pages (``Lecture 1 - slide 3 / 12`` and ``... 4 / 12`` are the same
    line) - are dropped before layout; the rule needs at least two pages. PDFs carry
    no speaker notes (``notes`` is ``None``). Each raster image on a page becomes a
    ``SlideImage`` (``_page_images``), listed once per slide and once in ``assets``;
    like ``parse_pptx`` this reports every image and leaves the size and recurring
    rules to ``ingest_slides``. Raises ``DeckParseError`` for a file pypdf cannot
    read or decrypt with the empty password; lets ``FileNotFoundError`` through.
    """
    if not path.is_file():
        raise FileNotFoundError(path)
    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted and reader.decrypt("") == PasswordType.NOT_DECRYPTED:
            raise DeckParseError(path, ValueError("password-protected PDF"))
        pages = [
            (
                _page_spans(page),
                _page_images(page),
                float(page.mediabox.width),
                float(page.mediabox.height),
            )
            for page in reader.pages
        ]
    except (PyPdfError, DependencyError) as exc:
        raise DeckParseError(path, exc) from exc

    recurring = _boilerplate([spans for spans, _, _, _ in pages])
    slides: list[Slide] = []
    assets: dict[str, SlideImage] = {}
    for number, (spans, images, width, height) in enumerate(pages, start=1):
        layout = layout_page(
            (span for span in spans if _normalised(span) not in recurring),
            page_width=width,
            page_height=height,
        )
        image_ids: list[str] = []
        for asset in images:
            assets.setdefault(asset.id, asset)
            if asset.id not in image_ids:
                image_ids.append(asset.id)
        slides.append(
            Slide(
                number=number,
                title=layout.title,
                blocks=layout.blocks,
                notes=None,
                image_ids=tuple(image_ids),
            )
        )
    return Deck(source=path.as_posix(), slides=tuple(slides), assets=tuple(assets.values()))


# --- image rules (pure, format-agnostic) ---------------------------------------------
#
# Both live here rather than in the parsers so they are written and tested once and
# behave identically for both formats - the same structure as ``merge_sentences``
# sitting behind ``ingest_captions``. They take and return ``Deck`` values (the models
# are frozen) so P2-04's command and Phase 4 tests can call them on hand-built decks.

_RECURRING_MIN_SLIDES = 3


def _with_image_ids(deck: Deck, keep: Callable[[str], bool]) -> tuple[Slide, ...]:
    return tuple(
        slide.model_copy(update={"image_ids": tuple(i for i in slide.image_ids if keep(i))})
        for slide in deck.slides
    )


def drop_small_images(deck: Deck, *, min_px: int) -> Deck:
    """Remove every asset narrower or shorter than ``min_px`` from ``assets`` and from
    every slide's ``image_ids``.

    Bullet glyphs and rules are 8–20 px; real figures are hundreds. Either dimension
    counts, so a 200 × 8 rule is a decoration too.
    """
    small = {a.id for a in deck.assets if a.width < min_px or a.height < min_px}
    if not small:
        return deck
    return deck.model_copy(
        update={
            "assets": tuple(a for a in deck.assets if a.id not in small),
            "slides": _with_image_ids(deck, lambda i: i not in small),
            "recurring_image_ids": tuple(i for i in deck.recurring_image_ids if i not in small),
        }
    )


def set_aside_recurring_images(deck: Deck) -> Deck:
    """Move the logo out of the way: an id on more than half of the slides of a deck
    with at least three leaves every ``image_ids``, stays in ``assets``, and is appended
    to ``recurring_image_ids`` in first-seen order.

    Phase 5 iterates ``image_ids`` and never sees it; a debugger (P2-04's ``slides``
    command) can still see what was set aside. "More than half of ≥ 3" because a
    two-slide deck cannot tell a logo from a figure shown twice. Idempotent.
    """
    if len(deck.slides) < _RECURRING_MIN_SLIDES:
        return deck
    counts = Counter(i for slide in deck.slides for i in slide.image_ids)
    recurring = [i for i, n in counts.items() if n > len(deck.slides) / 2]  # first-seen order
    if not recurring:
        return deck
    return deck.model_copy(
        update={
            "slides": _with_image_ids(deck, lambda i: i not in recurring),
            "recurring_image_ids": deck.recurring_image_ids + tuple(recurring),
        }
    )


# --- the composed stage ------------------------------------------------------------

_PARSERS: dict[str, Callable[[Path], Deck]] = {".pptx": parse_pptx, ".pdf": parse_pdf}


def ingest_slides(path: Path, *, min_px: int = 32) -> Deck:
    """Plan §3 stage 2 end to end: parse ``path`` by suffix, then apply the image rules.

    ``.pptx`` and ``.pdf``. After parsing, ``drop_small_images(min_px)`` then
    ``set_aside_recurring_images`` (a tiny logo on every slide is simply dropped).
    ``min_px`` is a knob because a course whose decks use tiny icons meaningfully can
    lower it. Raises ``ValueError`` for any other suffix, ``FileNotFoundError`` for a
    missing file, ``DeckParseError`` for a file that is not a deck.
    """
    suffix = path.suffix.lower()
    parser = _PARSERS.get(suffix)
    if parser is None:
        raise ValueError(f"unsupported deck format: {suffix!r} (expected .pptx or .pdf)")
    return set_aside_recurring_images(drop_small_images(parser(path), min_px=min_px))
